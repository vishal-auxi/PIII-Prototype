from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches


def create_bar_chart(details):
    if len(details["categories"]) != len(details["values"]):
        return False, "Couldn't create the bar chart as the number of categories and values are not equal"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    chart_data = CategoryChartData()

    try:
        values = (float(i) for i in details["values"])
    except ValueError as err:
        print(f'Cannot extract values from list: {err}')
        return False, "Couldn't create the bar chart due to invalid numeric values"

    chart_data.add_series('Series 1', values)
    chart_data.categories = details["categories"]

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    if 'title' in details:
        slide.placeholders[0].text = details['title'].title()
        # chart.chart_title.text_frame.text = details['title'].title()

    prs.save('Presentation - Bar Chart.pptx')

    return True, "Created the Bar Chart"
