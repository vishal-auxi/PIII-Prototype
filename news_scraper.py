import requests
import os
from pprint import pprint
import urllib
import urllib.request as urllib2
import os
from pptx import Presentation

NY_BASE = '../base_pres.pptx'


def news_scraper(topic):
    apikey = "OuCSULEj80JbINMUCsuKHpoMOcuilkIi"
    title = []
    body = []
    img_list = []
    query = topic
    begin_date = "20190701"  # YYYYMMDD
    # filter_query = "\"body:(\"Trump\") AND glocations:(\"WASHINGTON\")\""
    page = "0"  # <0-100>
    sort = "relevance"  # newest, oldest
    query_url = f"https://api.nytimes.com/svc/search/v2/articlesearch.json?" \
                f"q={query}" \
                f"&api-key={apikey}" \
                f"&begin_date={begin_date}" \
                f"&page={page}" \
                f"&sort={sort}"

    r = requests.get(query_url)
    response = r.json()
    for i in range(0, len(response['response']['docs'])):
        title.append(response['response']['docs'][i]['headline']['main'])
        body.append(response['response']['docs'][i]['lead_paragraph'])
        url = str('http://static01.nyt.com/' + response['response']['docs'][i]['multimedia'][0]['url'])
        urllib.request.urlretrieve(url, os.path.basename("/images/" + url))
        img_list.append("../images/" + url.split("/")[-1])
    return (title, body, img_list)


def alexa_pres(num_slides: int, headers_list: list, text_list: list, img_list: list):
    pres = Presentation('../base_pres.pptx')
    slides = pres.slides
    slide = slides[0]
    shapes = slide.shapes
    title = shapes.title
    title.text = headers_list[0]
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == 1:
                shape.text = text_list[0]
    if img_list[0] is not '':
        picture_placeholder = slide.placeholders[13]
        placeholder_picture = picture_placeholder.insert_picture(img_list[0])

    for x in range(1, num_slides - 1):
        slide = pres.slides.add_slide(pres.slide_layouts[3])
        shapes = slide.shapes
        title = shapes.title
        title.text = headers_list[x]
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.idx == 0:
                    print('skip title')
                if phf.idx == 1:
                    shape.text = text_list[x]
        if img_list[x] is not '':
            picture_placeholder = slide.placeholders[13]
            placeholder_picture = picture_placeholder.insert_picture(img_list[x])
    pres.save('Final - Presentation ny.pptx')
    # s3_client = boto3.client('s3', region_name='us-west-2')  # Change as appropriate
    # s3_client.upload_file('/tmp/yourpresentation.pptx', 'presslides', 'Presentation/yourpresentation.pptx')
    return pres


if not os.path.exists("images/"):
    os.mkdir("images/")
os.chdir("images/")

# (title, body, img_list) = news_scraper("politics")
#
# alexa_pres(10, title, body, img_list)


def create_ppt_ny(ppt_subject, count):
    try:
        (title, body, img_list) = news_scraper(ppt_subject)
        alexa_pres(count, title, body, img_list)
        return True, 'Created the presentation from New York Times'
    except Exception as err:
        print(f'Cannot create ny presentation: {err}')

    return False, 'Cannot create presentation from New York Times'
