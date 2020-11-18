from pptx import Presentation, exc
import json

details_path = './details.json'


def create_team_slide(people, base=None):
    count = len(people)

    if count > 4:
        return False, "Couldn't create the team slide as there were too many people"

    with open(details_path) as json_file:
        data = json.load(json_file)

    # prs = Presentation(base)

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        # print(f'No presentation file: {err}')
        prs = Presentation(base)

    details = [data["user"]]
    for i in people:
        found = False
        for j in data["people"]:
            if i.lower() in j["name"].lower():
                found = True
                details.append(j)
        if not found:
            return False, "Couldn't create the team slide as details for " + i + " were not available"

    ph_team_name = 10
    if count == 0:
        slide_layout = prs.slide_layouts[0]
        ph_team_name = 10
        ph_images = [11]
        ph_names = [16]
        ph_roles = [21]
        ph_about = [22]
    elif count == 1:
        slide_layout = prs.slide_layouts[1]
        ph_team_name = 10
        ph_images = [11, 22]
        ph_names = [16, 23]
        ph_roles = [21, 24]
        ph_about = [25, 26]
    elif count == 2:
        slide_layout = prs.slide_layouts[2]
        ph_team_name = 10
        ph_images = [11, 22, 25]
        ph_names = [16, 23, 26]
        ph_roles = [21, 24, 27]
        ph_about = [28, 29, 30]
    elif count == 3:
        slide_layout = prs.slide_layouts[3]
        ph_team_name = 10
        ph_images = [11, 22, 25, 28]
        ph_names = [16, 23, 26, 29]
        ph_roles = [21, 24, 27, 30]
        ph_about = [31, 32, 33, 34]
    elif count == 4:
        slide_layout = prs.slide_layouts[4]
        ph_names = [16, 17, 18, 19, 20]
        ph_roles = [21, 22, 23, 24, 25]
        ph_images = [11, 12, 13, 14, 15]
        ph_about = [31, 32, 33, 34, 35]

    slide = prs.slides.add_slide(slide_layout)

    # for shape in slide.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))

    phs = slide.placeholders
    phs[ph_team_name].text = data["team"]
    for i in range(len(details)):
        phs[ph_names[i]].text = details[i]["name"]
        phs[ph_roles[i]].text = details[i]["role"]
        phs[ph_images[i]].insert_picture(details[i]["img"])
        phs[ph_about[i]].text = details[i]["about"]

    prs.save('Final - Presentation.pptx')
    return True, "Created the team slide"

# def create_team_slide(people):
#     count = len(people)
#
#     if count > 4:
#         return False, "Couldn't create the team slide as there were too many people"
#
#     with open('../../misc/Details/details.json') as json_file:
#         data = json.load(json_file)
#
#     prs = Presentation('../../misc/Team_Slide_Custom.pptx')
#
#     # slide_layout = prs.slide_layouts[1]
#     # print(slide_layout)
#     # slide = prs.slides.add_slide(slide_layout)
#
#     slides = prs.slides
#     slide = slides[0]
#
#     for shape in slide.placeholders:
#         print('%d %s' % (shape.placeholder_format.idx, shape.name))
#
#     details = [data["user"]]
#     for i in people:
#         found = False
#         for j in data["people"]:
#             if i.lower() in j["name"].lower():
#                 found = True
#                 details.append(j)
#         if not found:
#             return False, "Couldn't create the team slide as details for " + i + " were not available"
#
#     ph_team_name = 10
#     if count == 0:
#         ph_names = [18]
#         ph_roles = [23]
#         ph_images = [13]
#         ph_remove_list = [11, 12, 14, 15, 16, 17, 19, 20, 21, 22, 24, 25]
#     elif count == 1:
#         ph_names = [17, 19]
#         ph_roles = [22, 24]
#         ph_images = [12, 14]
#         ph_remove_list = [11, 13, 15, 16, 18, 20, 21, 23, 25]
#     elif count == 2:
#         ph_names = [17, 18, 19]
#         ph_roles = [22, 23, 24]
#         ph_images = [12, 13, 14]
#         ph_remove_list = [11, 15, 16, 20, 21, 25]
#     elif count == 3:
#         ph_names = [16, 17, 18, 19]
#         ph_roles = [21, 22, 23, 24]
#         ph_images = [11, 12, 13, 14]
#         ph_remove_list = [15, 20, 25]
#     elif count == 4:
#         ph_names = [16, 17, 18, 19, 20]
#         ph_roles = [21, 22, 23, 24, 25]
#         ph_images = [11, 12, 13, 14, 15]
#         ph_remove_list = []
#
#     phs = slide.placeholders
#     phs[ph_team_name].text = data["team"]
#     for i in range(len(details)):
#         phs[ph_names[i]].text = details[i]["name"]
#         phs[ph_roles[i]].text = details[i]["role"]
#         phs[ph_images[i]].insert_picture(details[i]["img"])
#
#     for i in ph_remove_list:
#         placeholder = phs[i]
#         sp = placeholder._sp
#         sp.getparent().remove(sp)
#
#     # ph_names = [24, 25, 26]
#     # ph_roles = [28, 29, 30]
#     # ph_images = [21, 22, 23]
#     #
#     # phs = slide.placeholders
#     # phs[ph_team_name].text = team
#     # for i in range(count):
#     #     phs[ph_names[i]].text = people[i]["name"]
#     #     phs[ph_roles[i]].text = people[i]["role"]
#     #     phs[ph_images[i]].insert_picture(people[i]["img"])
#
#     prs.save("Presentation - Team Slide.pptx")
#     return True, "Created the team slide"
