from pptx import Presentation, exc
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches


def create_line_chart(details, base=None):

    if len(details["Label"]) != len(details["Value"]):
        return False, "Couldn't create the line chart as the number of labels and values are not equal"

    # prs = Presentation()
    # slide = prs.slides.add_slide(prs.slide_layouts[5])

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        # print(f'No presentation file: {err}')
        prs = Presentation(base)
    slide = prs.slides.add_slide(prs.slide_layouts[10])

    chart_data = CategoryChartData()
    # chart_data = ChartData()

    try:
        values = (float(i) for i in details["Value"])
    except ValueError as err:
        print(f'Cannot extract values from list: {err}')
        return False, "Couldn't create the line chart due to invalid numeric values"

    chart_data.add_series('Series 1', values)
    chart_data.categories = details["Label"]

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[0].smooth = True

    if 'title' in details:
        slide.placeholders[0].text = details['title'].title()
        # chart.chart_title.text_frame.text = details['title'].title()

    prs.save('Final - Presentation.pptx')

    return True, "Created the Line Chart"
