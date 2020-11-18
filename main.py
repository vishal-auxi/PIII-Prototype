from flask import Flask, request, jsonify
from pptx import Presentation, exc
import wikipedia

from Summarizer import summarize_lsa, sentence_count
from team_slide import create_team_slide
from pie_chart import create_pie_chart
from bar_chart import create_bar_chart
from line_chart import create_line_chart

from news_scraper import create_ppt_ny

import json

import os

BASE_PRESENTATION = "./Team_Slide_Custom.pptx"

app = Flask(__name__)


def get_summarized_sections(ppt_subject):
    pages_search = wikipedia.search(ppt_subject, results=5, suggestion=False)
    page_suggest = wikipedia.suggest(ppt_subject)

    if not pages_search and not page_suggest:
        return False

    page_match = ppt_subject
    if pages_search:
        print(pages_search)
        for page in pages_search:
            print("in loop for i = " + page)
            try:
                wikipedia.page(page).content
                page_match = page
                break
            except wikipedia.DisambiguationError as e:
                continue
            except Exception as e:
                continue
        # page_match = pages_search[0]
    elif page_suggest:
        page_match = page_suggest

    page = wikipedia.page(page_match)

    title = page.title
    url = page.url
    summary = page.summary
    images = page.images
    categories = page.categories
    content = page.content
    links = page.links
    sections = page.sections
    references = page.references

    if not sections:
        return False

    sections_summarized = {}
    for section in sections:
        if section.lower() == "see also":
            break
        sections_summarized[section] = summarize_lsa(page.section(section))

    return title, sections_summarized, url


def create_ppt(ppt_subject):
    print("in create_ppt")

    result = get_summarized_sections(ppt_subject)
    if not result:
        return (False,)

    page_title = result[0]
    sections = result[1]
    url = result[2]

    # for section in sections.items():
    #     print(section[0] + ":")
    #     print(section[1])
    #     print("\n")

    # prs = Presentation()
    # title_slide_layout = prs.slide_layouts[0]

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        # print(f'No presentation file: {err}')
        prs = Presentation(BASE_PRESENTATION)
    title_slide_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = page_title
    subtitle.text = "Created by Presentation Bot \n (Sourced from {url})".format(url=url)

    for section in sections.items():
        if not section[1]:
            continue

        # slide_layout = prs.slide_layouts[1]
        slide_layout = prs.slide_layouts[6]

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = section[0]
        subtitle_text = ""
        for i in section[1]:
            subtitle_text = subtitle_text + "\n" + i
        # subtitle.text = section[1][0]
        subtitle.text = subtitle_text

    # prs.save("Presentation - " + page_title + ".pptx")
    prs.save('Final - Presentation.pptx')

    return True, page_title


def create_ppt_with_count(ppt_subject, count):
    print("in create_ppt_with_count")

    result = get_summarized_sections(ppt_subject)
    if not result:
        return (False,)

    page_title = result[0]
    sections = result[1]
    url = result[2]

    # prs = Presentation()
    # title_slide_layout = prs.slide_layouts[0]

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        # print(f'No presentation file: {err}')
        prs = Presentation(BASE_PRESENTATION)
    title_slide_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = page_title
    subtitle.text = "Created by Presentation Bot \n (Sourced from {url})".format(url=url)

    total = 0
    for section in sections.items():
        if total >= count:
            break

        if not section[1]:
            continue

        # slide_layout = prs.slide_layouts[1]
        slide_layout = prs.slide_layouts[6]

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = section[0]
        subtitle_text = ""
        for i in section[1]:
            subtitle_text = subtitle_text + "\n" + i
        # subtitle.text = section[1][0]
        subtitle.text = subtitle_text

        total += 1

    # prs.save("Presentation - " + page_title + ".pptx")
    prs.save('Final - Presentation.pptx')

    return True, page_title, total


def create_charts(req_body):
    details = req_body['details']

    success = True
    response = ''

    if 'create_pie_chart' in details and details['create_pie_chart'] and details['pie_data']:
        pie_response = create_pie_chart(details['pie_data'], BASE_PRESENTATION)
        success = success and pie_response[0]
        response = "{pr}".format(response=response, pr=pie_response[1])

    if 'create_bar_chart' in details and details['create_bar_chart'] and details['bar_data']:
        bar_response = create_bar_chart(details['bar_data'], BASE_PRESENTATION)
        success = success and bar_response[0]
        response = "{response}. {br}".format(response=response, br=bar_response[1])

    if 'create_line_chart' in details and details['create_line_chart'] and details['line_data']:
        line_response = create_line_chart(details['line_data'], BASE_PRESENTATION)
        success = success and line_response[0]
        response = "{response}. {lr}".format(response=response, lr=line_response[1])

    if success:
        return True, response
    else:
        return False, response


@app.route("/", methods=['POST'])
def ppt_request_handle():
    print("in ppt_request_handle")
    req = request.json
    print(req)

    with open("req_sample.json", "w") as outfile:
        json.dump(req, outfile)

    if req["req"] == "create_ppt" and req["source"].lower() == "wikipedia":
        result = create_ppt(req["title"])
        if result[0]:
            response = {
                "message": "Request successful",
                "title": result[1]
            }
            return jsonify(response), 200

    elif req["req"] == "create_ppt_count" and req["source"].lower() == "wikipedia":
        result = create_ppt_with_count(req["title"], req["count"])
        if result[0]:
            response = {
                "message": "Request successful",
                "title": result[1],
                "count": result[2],
            }
            return jsonify(response), 200

    elif req["req"] == "create_ppt_count" and req["source"].lower() == "nytimes":
        result = create_ppt_ny(req["title"], req["count"])
        if result[0]:
            response = {
                "message": "Request successful",
                "title": result[1],
                "count": result[2],
            }
            return jsonify(response), 200

    elif req["req"] == "create_team_slide":
        result = create_team_slide(req["people"])
        response = {
            "message": result[1]
        }
        if result[0]:
            return jsonify(response), 200
        else:
            return jsonify(response), 400

    elif req["req"] == "create_pie_chart":
        result = create_pie_chart(req, BASE_PRESENTATION)
        response = {
            "message": result[1]
        }
        if result[0]:
            return jsonify(response), 200
        else:
            return jsonify(response), 400

    elif req["req"] == "create_bar_chart":
        result = create_bar_chart(req, BASE_PRESENTATION)
        response = {
            "message": result[1]
        }
        if result[0]:
            return jsonify(response), 200
        else:
            return jsonify(response), 400

    elif req["req"] == "create_line_chart":
        result = create_line_chart(req, BASE_PRESENTATION)
        response = {
            "message": result[1]
        }
        if result[0]:
            return jsonify(response), 200
        else:
            return jsonify(response), 400

    elif req["req"] == "create_chart" or req["req"] == "create_charts":
        result = create_charts(req)
        response = {
            "message": result[1]
        }
        if result[0]:
            return jsonify(response), 200
        else:
            return jsonify(response), 400

    elif req["req"] == "create_org_chart":
        result = True, 'Got the request'
        response = {
            "message": result[1]
        }
        if result[0]:
            return jsonify(response), 200
        else:
            return jsonify(response), 400

    return f"Cannot process request", 400


# 'Main' function to run
if __name__ == '__main__':
    # app.run(debug=True)  # run server in debug mode

    os.chdir("/Users/vishal/Documents/Auxi/PIII-Prototype/PIII-Prototype-First")
    with open('./1.json') as json_file:
        data = json.load(json_file)

    print(data)
    create_ppt_ny(data['title'], data['count'])
    # create_ppt(data['title'])

    # try:
    #     prs = Presentation('./Final - Presentation.pptx')
    # except exc.PackageNotFoundError as err:
    #     # print(f'No presentation file: {err}')
    #     prs = Presentation('../../misc/Team_Slide_Custom.pptx')

    # create_team_slide(["chris", "john", "mike", "steve"], BASE_PRESENTATION)
    # res = create_pie_chart({
    #     'req': 'create_pie_chart',
    #     'categories': ['USA', 'Canada', 'Mexico'],
    #     'percentages': ['30', '30', '40']}, BASE_PRESENTATION
    # )
    #
    # res = create_bar_chart({
    #     'req': 'create_bar_chart',
    #     'categories': ['USA', 'Canada', 'Mexico'],
    #     'values': ['81', '45', '54']}, BASE_PRESENTATION
    # )
    #
    # res = create_line_chart({
    #     'req': 'create_line_chart',
    #     "Label": ["India", "Bangladesh"],
    #     "Value": [20, 30]}, BASE_PRESENTATION
    # )

    # req = {
    #     "req": "create_chart",
    #     "details": {
    #         "create_pie_chart": True,
    #         "pie_data": {
    #             "categories": ["Russia", "India", "Bangladesh"],
    #             "percentages": [30.0, 30.0, 40.0]
    #         },
    #         "create_bar_chart": True,
    #         "bar_data": {
    #             "categories": ["Russia", "Russia", "Russia"],
    #             "values": [20, 20]
    #         },
    #         "create_line_chart": True,
    #         "line_data": {
    #             "Label": ["India", "Bangladesh"],
    #             "Value": [20, 30]
    #         }
    #     }
    # }
    #
    # res = create_charts(req)
    #
    # print(res)
