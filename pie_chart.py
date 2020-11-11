from pptx import Presentation, exc
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches


def create_pie_chart(details, base=None):

    if len(details["categories"]) != len(details["percentages"]):
        return False, "Couldn't create the pie chart as the number of categories and percentage values are not equal"

    # prs = Presentation()
    # slide = prs.slides.add_slide(prs.slide_layouts[5])

    try:
        prs = Presentation('./Final - Presentation.pptx')
    except exc.PackageNotFoundError as err:
        # print(f'No presentation file: {err}')
        prs = Presentation(base)
    slide = prs.slides.add_slide(prs.slide_layouts[10])

    chart_data = CategoryChartData()

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

    chart_data = ChartData()
    # chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
    # chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))

    try:
        percent_floats = [float(i) for i in details["percentages"]]
    except ValueError as err:
        print(f'Cannot extract number from percentages: {err}')
        return False, "Couldn't create the pie chart due to invalid percentage values"

    if sum(percent_floats) != 100.0 and sum(percent_floats) != 1.0:
        return False, "Couldn't create the pie chart as percentages don't add up to 100"

    if sum(percent_floats) == 100.0:
        temp = percent_floats
        percent_floats = [i / 100.0 for i in temp]

    chart_data.add_series('Series 1', percent_floats)

    chart_data.categories = details["categories"]

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    if 'title' in details:
        slide.placeholders[0].text = details['title'].title()
        # chart.chart_title.text_frame.text = details['title'].title()

    prs.save('Final - Presentation.pptx')

    return True, "Created the Pie Chart"
